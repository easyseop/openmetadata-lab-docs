#!/usr/bin/env python3
# CODEX-05 전략 발표 (v3) — 검사기를 정보구조의 중심으로: 라이프사이클 + 7기능표 + 현황표
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x1F,0x2A,0x44); INK=RGBColor(0x14,0x1B,0x2B); BLUE=RGBColor(0x2E,0x6C,0xB1); TEAL=RGBColor(0x0E,0x8F,0x8F)
AMBER=RGBColor(0xB8,0x79,0x1F); RED=RGBColor(0xB0,0x3A,0x2E); GREEN=RGBColor(0x2E,0x8B,0x57)
GRAY=RGBColor(0x55,0x5B,0x66); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x22,0x27,0x30)
ROW1=RGBColor(0xFF,0xFF,0xFF); ROW2=RGBColor(0xEE,0xF2,0xF7)
KFONT="맑은 고딕"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _set(r,s,c,b=False,f=KFONT): r.font.size=Pt(s); r.font.color.rgb=c; r.font.bold=b; r.font.name=f
def box(sl,l,t,w,h):
    tb=sl.shapes.add_textbox(l,t,w,h); tb.text_frame.word_wrap=True; return tb.text_frame
def rect(sl,l,t,w,h,fill):
    sp=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def footer(sl,n):
    tf=box(sl,Inches(0.4),Inches(7.06),Inches(10),Inches(0.33))
    r=tf.paragraphs[0].add_run(); r.text="OpenMetadata 커스텀 버전업 검증 전략 · CODEX-05(정본) 기반"; _set(r,9,GRAY)
    tf2=box(sl,Inches(12.4),Inches(7.06),Inches(0.7),Inches(0.33)); p=tf2.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    r=p.add_run(); r.text=str(n); _set(r,9,GRAY)
def head(title,kicker,accent):
    s=prs.slides.add_slide(BLANK)
    rect(s,0,0,SW,Inches(1.15),NAVY); rect(s,0,Inches(1.15),SW,Inches(0.06),accent)
    kf=box(s,Inches(0.5),Inches(0.18),Inches(12.2),Inches(0.35))
    r=kf.paragraphs[0].add_run(); r.text=kicker; _set(r,12,accent,True)
    tf=box(s,Inches(0.5),Inches(0.5),Inches(12.3),Inches(0.62))
    r=tf.paragraphs[0].add_run(); r.text=title; _set(r,25,WHITE,True)
    return s
def content(title,kicker,bullets,accent=BLUE,n=0):
    s=head(title,kicker,accent)
    bf=box(s,Inches(0.6),Inches(1.45),Inches(12.1),Inches(5.4)); first=True
    for it in bullets:
        text,lvl,col=(it if isinstance(it,tuple) else (it,0,DARK))
        p=bf.paragraphs[0] if first else bf.add_paragraph(); first=False
        p.space_after=Pt(6); p.space_before=Pt(2)
        if lvl==0:
            r=p.add_run(); r.text="■ "; _set(r,15,accent,True); r=p.add_run(); r.text=text; _set(r,15,col,True)
        elif lvl==1:
            p.level=1; r=p.add_run(); r.text="– "; _set(r,13,GRAY); r=p.add_run(); r.text=text; _set(r,13,col)
        else:
            p.level=2; r=p.add_run(); r.text="· "+text; _set(r,12,GRAY)
    footer(s,n); return s
def table_slide(title,kicker,headers,rows,widths,accent,n,note=None):
    s=head(title,kicker,accent)
    nrows=len(rows)+1; ncols=len(headers)
    gf=s.shapes.add_table(nrows,ncols,Inches(0.6),Inches(1.5),Inches(12.1),Inches(0.4*nrows))
    tb=gf.table; tb.first_row=False; tb.horz_banding=False
    for j,w in enumerate(widths): tb.columns[j].width=Inches(w)
    for j,h in enumerate(headers):
        c=tb.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY
        c.vertical_anchor=MSO_ANCHOR.MIDDLE; c.margin_top=Pt(3); c.margin_bottom=Pt(3)
        p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=h; _set(r,11,WHITE,True)
    for i,row in enumerate(rows,start=1):
        for j,val in enumerate(row):
            txt,col=(val if isinstance(val,tuple) else (val,DARK))
            c=tb.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=ROW1 if i%2 else ROW2
            c.vertical_anchor=MSO_ANCHOR.MIDDLE; c.margin_top=Pt(2); c.margin_bottom=Pt(2); c.margin_left=Pt(6)
            p=c.text_frame.paragraphs[0]; p.word_wrap=True; r=p.add_run(); r.text=txt
            _set(r,10.5,col, j==0)
    if note:
        nf=box(s,Inches(0.6),Inches(6.7),Inches(12.1),Inches(0.5))
        r=nf.paragraphs[0].add_run(); r.text=note; _set(r,10,GRAY)
    footer(s,n); return s

# ══════════════════════════════════════════════════════════════
# 슬라이드 정의 — HTML 덱(CODEX-05-전략발표.html, 16장)과 동일 구성
# ══════════════════════════════════════════════════════════════

# 1 표지
s=prs.slides.add_slide(BLANK); rect(s,0,0,SW,SH,INK); rect(s,0,Inches(4.05),SW,Inches(0.08),TEAL)
tf=box(s,Inches(0.9),Inches(2.0),Inches(11.6),Inches(2.0))
r=tf.paragraphs[0].add_run(); r.text="커스텀 버전업 검증 전략"; _set(r,42,WHITE,True)
p=tf.add_paragraph(); r=p.add_run()
r.text="규칙 기반 검사기에서 무엇을 LLM 체제로 가져갈 것인가"; _set(r,20,RGBColor(0xBF,0xD3,0xE6))
sf=box(s,Inches(0.9),Inches(4.35),Inches(11.6),Inches(1.6))
for i,t in enumerate([
  "1.13.0→1.13.1 통과 · 1.13.1→1.13.2에서 드러난 한계 · 그래서 정한 방향",
  "관리 파일 18종 1차 검토 결과 포함 · 정본 CODEX-05 기반","2026-08-12"]):
    p=sf.paragraphs[0] if i==0 else sf.add_paragraph(); r=p.add_run(); r.text=t; _set(r,14,RGBColor(0x9F,0xB2,0xC8))

# 2 배경
content("배경과 처음의 생각","왜 시작했나",
 [("목표 — 우리가 손댄 부분(DB 커넥터·자체 화면·한글 입력 등)이 공식 버전업 뒤에도 안 깨지고 살아 있는지 확인한다",0,DARK),
  ("전제 — AI를 쓰는 것 자체는 옳다. 다만 행내에서 운영하는 서비스라 최대한 보수적으로 가야 한다고 봤다",0,DARK),
  ("그래서 AI 개입을 최소화하고, 사람·AI 주관을 뺀 채 정해진 규칙(파일 위치·커밋 이력)만 보는 자동 검사로 전부 걸러내는 전략을 먼저 검토했다",1,DARK),
  ("효과 — 같은 코드면 항상 같은 결과가 나오는 ‘재현 가능한 자동 판정’",0,GREEN)],BLUE,2)

# 3 초기 구현
content("처음 만든 것 — 규칙 기반 자동 검사기","어떻게 만들었나 (초기)",
 [("파일 위치 — ‘이 기능은 이 파일들에’를 목록에 적고 대조",0,DARK),
  ("커밋 이력 — 커밋마다 꼬리표(Customization-ID) + 선언 범위 밖 변경 확인",0,DARK),
  ("안전장치 — 애매하면 통과가 아닌 ‘막음’(fail-closed)",0,DARK),
  ("못박기 — 검사한 코드를 특정 커밋·tree에 고정",0,DARK),
  ("예시) BANK-OM-001(InstanceCode) 하나가 만드는 관리 문서 2종",0,TEAL),
  ("manifests/BANK-OM-001.yaml — 경로 목록(‘경로 기반’의 실체). allowed 약 50개 / required는 핵심 앵커",1,DARK),
  ("contracts.yaml — 불변식(codeGroup/codeValue가 CRUD·재색인 뒤에도 동일 조회) + 그걸 지키는 테스트",1,DARK),
  ("여기 보인 건 2종일 뿐 — registry·shared-path-owners·정책 3종·계획 3종·못박기·증거 3종까지 합쳐 커스텀 11개 기준 약 25개의 관리 파일이 함께 돈다 (→ 6페이지)",2,GRAY)],BLUE,3)

# 4 업그레이드 테스트
content("업그레이드 테스트","직접 돌려봄",
 [("1.13.0 → 1.13.1  성공 — 변화가 작아 충돌 0·자동 병합. 규칙 기반으로 깔끔히 통과",0,GREEN),
  ("1.13.1 → 1.13.2  여기서 막힘 — 아래 ①→② 순서",0,RED),
  ("① 병합(먼저) — 충돌 다수 + Sybase 코드가 다른 파일로 이사(리로케이션). 커스텀을 새 위치에 수작업 재적용하고 관리 파일 ~25개·커밋 이력도 함께 손봐야 했다",1,DARK),
  ("② 실행·검증(그 뒤) — 병합을 수작업 해결한 뒤 이미지 빌드 → 데이터 이관 → 서버 기동 → API 확인. InstanceCode는 실제 서버에서 CRUD·검색까지 동작 (PASS)",1,DARK),
  ("업그레이드 테스트를 하며 겪은 어려움 — 리로케이션은 드러난 계기 중 하나일 뿐이었다",0,RED),
  ("① 코드 이동 — 코드가 어디로 옮겨갔는지 규칙으론 못 따라간다 (옛 경로만 보고 “사라졌다”고 오판)",1,DARK),
  ("② 반쪽 병합 — 검사기는 ‘그 파일이 공식과 다른가’로만 판정한다. 어떻게 다른지는 안 본다 → 절반만 들어가도 통과된다",1,DARK),
  ("③ 등록할 경로가 많다 — Sybase 하나에 18개, 그중 11개가 자동 생성 파일이라 다시 만들 때마다 통째로 바뀐다",1,DARK),
  ("④ 관리 부담 — 그때마다 관리 파일 ~25개와 커밋 이력을 사람이 갱신해야 한다",1,DARK)],RED,4)

# 5 못 보는 것들
content("규칙 검사만으로는 못 보는 것들","드러난 한계 ① — 규칙 검사의 사각",
 [("리로케이션 — 코드가 다른 파일로 이사 가면 옛 경로만 보고 ‘사라졌다’고 오판한다",0,RED),
  ("의미(반쪽 병합) — 검사기의 ‘생존’ 판정 기준은 “커스텀 파일이 공식과 다른가”다. 그런데 ‘다르다’ ≠ ‘기능이 온전하다’",0,RED),
  ("병합이 어긋나 커스텀 로직이 반쪽만 들어가도 ‘공식과 다름’은 성립 → survived(통과)로 새어나간다",1,DARK),
  ("DB 마이그레이션이 어긋나 서버가 아예 안 뜬 적도 있다 — 파일은 멀쩡한데 런타임이 깨진 경우. 경로만 보는 검사론 못 잡는다",1,DARK)],RED,5)

# 6 지속 운영 한계
content("지속 가능 운영의 한계 — 관리할 파일이 너무 많다","드러난 한계 ② — 지속 가능 운영",
 [("부담 — 무엇을 바꿨는지·어디까지 허용인지·누가 승인했는지를 전부 사람이 파일로 관리한다. 반복 버전업 때마다 이 목록들을 일일이 갱신해야 한다",0,RED),
  ("본질 — ‘검사 자체’보다 ‘검사를 유지하는 비용’이 더 크다. 지속 운영의 벽",0,RED),
  ("위험 — 목록이 늘수록 갱신 누락·중복 검사·피로로 인한 형식적 승인 위험이 커진다",1,DARK)],RED,6)

# 7 강점
content("강점 — LLM 혼자선 보장 못 하는 것들","왜 ‘룰베이스 검사기’인가",
 [("아래는 전부 LLM도 할 수 있는 일이다. 차이는 능력이 아니라 누가 증언하느냐에 있다 — LLM이 하면 LLM의 말이 되고, 검사기가 하면 증거가 된다",0,TEAL),
  ("허위 테스트 방지 — 기능이 되는지가 아니라 그 테스트가 진짜인지를 검사한다",0,GREEN),
  ("① 테스트 사전 정의 → ② 공식 버전에서 FAIL 이어야 정상 → ③ 커스텀 버전에서 PASS 여야 정상. ②가 PASS면 속 빈 테스트",1,DARK),
  ("LLM이 하면: “돌려봤고 실패했습니다”—사실인지 알 방법이 없다 / 검사기가 하면: 어느 커밋에서 어떻게 실패했는지 파일로 남는다",1,GRAY),
  ("한계 — 아무것도 안 보는 테스트는 걸러내지만, 덜 보는 테스트는 못 걸러낸다 (status_code==200만 봐도 통과)",2,AMBER),
  ("향후 개선 — 돌연변이 테스트: 커스텀을 통째로 빼는 대신 한 줄씩만 바꿔 그때도 테스트가 실패하는지 본다 → ‘덜 보는 테스트’까지 드러난다",2,AMBER),
  ("막고 보기(fail-closed) — 검사 하나라도 못 돌면 통과가 아니라 막음",0,GREEN),
  ("실제: runtime contracts 2 pass · 7 skipped → verdict: block  (실패가 0인데도 통과가 아니다)",1,DARK),
  ("못박기(candidate-lock) — 검사한 코드를 커밋에 고정해 이후 바뀌면 드러난다",0,GREEN),
  ("기계 대조(data-check) — 고정된 건 값이 아니라 판정 규칙. 값이 매번 달라지면 ‘넣은 값이 그대로 나오는지’(round-trip)로 대조",0,GREEN),
  ("결정론(재현성) — 이 항목만은 LLM이 구조적으로 못 한다. 같은 입력에도 답이 흔들려 배포 기준이 될 수 없다",0,BLUE)],GREEN,7)

# 8 테스트 결론
content("그래서 내린 판단","테스트 결론",
 [("처음의 규칙 기반 · AI 최소화 전략을 실제로 돌려본 결과 —",0,DARK),
  ("안 된 것 ① — 돌발 상황에는 유연하지 못했다. 규칙은 미리 정해둔 것만 보므로, 예상 못 한 변경이 오면 사람이 손으로 따라가야 한다",0,RED),
  ("안 된 것 ② — 유지 비용이 컸다. 그때마다 등록 경로와 관리 파일을 다시 맞춰야 한다",0,RED),
  ("잘 된 것 — 어려웠던 건 대처와 관리였지 판정이 아니었다. 정해진 형식 안에서 PASS/FAIL을 흔들림 없이 내는 일은 규칙 기반이 해냈다",0,GREEN),
  ("그 테스트가 진짜인지 · 안 돌린 걸 통과로 세지 않는지 · 검사한 코드가 그대로인지 · 같은 코드면 같은 판정인지 (7페이지의 강점들)",1,GRAY),
  ("→ 대처와 관리를 LLM에게 넘겨도 이 판정은 그대로 쓸 수 있다",1,DARK),
  ("판단 — 지속 가능한 운영이 어렵다. 그래서 기본 운영 구조를 LLM으로 옮기되, 통제는 놓지 않는다",0,TEAL),
  ("① 방식 통제 — 어떻게 일을 하는지를 지시(프롬프트·MD)로 통제한다 (순서·범위·남길 것)",0,TEAL),
  ("② 결과 통제 — 과정마다 정해진 형식으로 결과를 뱉게 강제 → 검사기가 판정. 자유서술로 “했다”고 적게 두지 않는다",0,GREEN),
  ("③ 최종 목표 — ①②를 ‘스킬’(정해진 작업 절차)로 굳혀 OpenMetadata 운영 하네스 템플릿을 갖추는 것 (아직 만들지 않았다)",0,BLUE)],TEAL,8)

# 9 8/5 접근
content("8/5 공유된 ‘LLM 기반 운영 방식’ (참고)","참고한 다른 접근",
 [("개념 — 규칙·지식 문서·자동화로 LLM이 버전 포팅을 빠르게 수행하는 방식",0,BLUE),
  ("“포팅 시간 크게 단축”은 그 방식의 자체 보고 수치 — 우리가 실행·검증한 건 아니다",1,GRAY),
  ("구조 — 변경·수정·테스트·결과 기록을 대부분 LLM이 스스로 수행하고, 그 결과도 스스로 문서에 적는다",0,DARK),
  ("우리 생각 — 검증이 아니라 방향 제안: 이 접근에 검사기로 ‘각 판단의 근거를 자동으로 뒷받침’하는 층을 더하면 속도(LLM)는 살리고 신뢰(근거)는 끌어올릴 수 있겠다. 경쟁이 아니라 보완",0,TEAL)],BLUE,9)

# 9 방향 전환
content("LLM이 일하고, 검사기가 근거를 붙인다","방향 전환",
 [("LLM — 판단·제안: 무엇을 바꿀지, 코드가 어디로 옮겨졌을지, 충돌 해결, 필요한 테스트",0,BLUE),
  ("검사기 — 뒷받침: 그 판단을 실제 코드·실행 증거로 확인하고, 위험하면 막는다",0,TEAL),
  ("구체적으로 — LLM의 주장을 검사기가 이렇게 확인한다",0,DARK),
  ("“이 파일 바꿨다” → 실제 커밋의 변경 내역을 직접 대조",1,DARK),
  ("“충돌 0건” → 임시로 실제 병합해 충돌 수를 센다 (안 해보고 적어도 티가 안 나는 주장)",1,DARK),
  ("“테스트 통과” → 실제 서버에서 다시 실행 (미실행은 통과가 아니다)",1,DARK),
  ("“승인함” → 승인을 결과에 묶어 입력이 바뀌면 자동 무효",1,DARK),
  ("핵심 규칙 — LLM이 ‘됐다’고 적어도 증거가 없으면 ‘미확인’으로 표시한다. 증거 없다고 곧장 ‘거짓’으로 몰지 않는다 → 확인·반박·미확인·사람검토 4단계",0,TEAL)],TEAL,10)

# 10 도입 제안 ① 검증 과정
content("커스텀 하나(BANK-OM ID)가 검증되는 과정","도입 제안 ① — 검증되는 과정",
 [("제안 — 하이브리드에서 LLM 작업의 ‘증거 층’으로 이 검사기를 도입한다 (자랑이 아니라 도입 대상)",0,TEAL),
  ("① 등록 — Manifest·Registry·Contract. 무엇을·어디를 바꿨고 어떤 검사로 지킬지 선언",1,DARK),
  ("② 커밋·변경범위 검사 — 선언한 범위만 건드렸는지, 공식 계보 위에 얹혔는지",1,DARK),
  ("③ 공용 파일의 ID별 코드 검사 — 공유 파일에서 이 ID 몫의 코드가 살아있는지",1,DARK),
  ("④ 전용 검사 실행 — 이 기능만의 검사(API/화면/소스). 속 빈 테스트면 patch-kill로 걸러낸다",1,DARK),
  ("⑤ commit·이미지와 결과 결속 — 바꿔치기·사후수정을 감지",1,DARK),
  ("⑥ 승인·증거 보관 — 권한자 승인 + 증거 보관. 승인은 결과에 결속",1,DARK),
  ("6페이지 한계와의 연결 — “관리 파일이 많다”는 부담을 그대로 두지 않는다. 이 산출물을 사람이 아니라 LLM이 생성·갱신하고, 검사기는 검증만 한다",0,TEAL)],TEAL,11)

# 11~13 관리 파일 18종
MF_A=[("customization-registry.yaml","우리가 손댄 기능 전체 명단","채택","LLM이 명단을 갱신, 사람은 확인만"),
 ("manifests/BANK-OM-00N.yaml (11개)","기능별 상세 카드","채택","리로케이션 대응이 여기로 — LLM이 새 위치를 찾아 경로를 고쳐 제안"),
 ("contracts.yaml","지켜야 할 약속 + 확인 테스트","채택","사람과 LLM의 경계선 — 불변식(무엇)은 사람이, 테스트(어떻게)는 LLM이"),
 ("shared-path-owners.yaml","공유 파일의 기능별 지분(37개)","채택","LLM이 초안 작성 (사람이 일일이 세던 일)"),
 ("source-diff-paths.txt","바뀐 파일 전체 목록(113)","채택","LLM이 읽고 ‘등록 누락’을 찾아 제안"),
 ("patch-kill-plan.yaml","빼보기 검사 계획(소스)","보류","LLM은 ‘이 테스트가 진짜다’를 스스로 주장할 수 없다"),
 ("runtime-patch-kill-plan.yaml","빼보기 검사 계획(실서버)","보류","공식 이미지로 한 번 돌리는 것으로 정리된다"),
 ("upgrade-rehearsal-matrix.yaml","버전 올림 예행연습 표","보류","LLM에게 병합을 맡길수록 격리가 중요해진다"),
 ("previous-run-test-review.yaml","직전 버전 테스트 기록 확인","신규","LLM이 가장 잘하는 일 — 직전 기록을 읽고 원인을 정리")]
MF_B=[("policies/sensitive-zones.yaml","손대면 위험한 구역","채택","LLM에게 주는 유일한 하드 제약. 이 파일 자체를 못 고친다"),
 ("policies/repository-layout.yaml","공식/우리 것 구분 지도","제외","폴더 구분 판단은 LLM이 한다"),
 ("policies/debt-thresholds.yaml","고친 양 상한선","제외","LLM은 필요하면 대규모로 고칠 수 있어야 한다"),
 ("change-intent-current.yaml","이번 작업의 허용/금지 범위","제외","LLM이 매니페스트에서 범위를 읽어 판단"),
 ("candidate-locks/ (lock+approval)","코드 봉인 + 승인 묶기","채택","LLM이 손댈 수 없는 영역. 판정과 승인은 못 바꾼다")]
MF_E=[("source-candidate-evidence.yaml","소스 검사 결과 기록","채택","LLM 주장과 대조할 원본. LLM은 이 파일을 만들지 않는다"),
 ("source-patch-kill-evidence.yaml","빼보기 결과 기록","보류","‘공식 버전에서 몇 개가 제대로 실패했나’가 기록될 자리"),
 ("ui-typecheck-baseline-evidence.yaml","화면 타입오류 기준선","보류","LLM이 옮겨 붙인 결과를 값싸게 검증"),
 ("previous-run-test-review-evidence.yaml","직전 기록 확인 결과","신규","LLM이 정리한 회고를 기록으로 고정")]
VC={"채택":GREEN,"제외":GRAY,"보류":AMBER,"신규":BLUE}
def mf_slide(title,kicker,rows,accent,n,note):
    return table_slide(title,kicker,["파일","용도","판정","LLM 체제에서의 역할"],
      [[f,u,(v,VC[v]),l] for f,u,v,l in rows],[3.5,2.5,0.8,5.3],accent,n,note)

mf_slide("LLM이 만들고, 사람이 확인한다  (관리 파일 9종)","도입 제안 ② — LLM이 만든다",MF_A,BLUE,12,
 "1차 검토 전체 18종 기준 — 채택 8 · 제외 3 · 보류 5 · 신규 2 제안   ※ 4페이지와 모순 아닌가: 버리는 것은 ‘경로가 안 맞으면 죽었다고 판정’하던 것. 남기는 것은 ‘어디서부터 찾을지의 출발점’ — 경로 목록이 없으면 LLM도 일을 못 한다")
mf_slide("사람이 정하고, LLM이 못 바꾼다  (관리 파일 5종)","도입 제안 ③ — LLM이 못 바꾼다",MF_B,TEAL,13,
 "다 덜어내도 이 둘은 남긴다 — ① 손대면 안 되는 구역(sensitive-zones) ② 검사한 코드와 승인을 못 바꾸게 묶는 것(candidate-locks)")
mf_slide("검사기가 자동으로 남긴다  (관리 파일 4종)","도입 제안 ④ — 검사기가 남긴다",MF_E,GREEN,14,
 "LLM은 “테스트 통과했습니다”라고 적는 데 비용이 0이다 → 검사기가 독립적으로 만든 이 기록만이 LLM 주장의 근거가 된다. 관리 부담은 원래 0인 영역")

# 14 그래서 빨라지나
table_slide("그래서 빨라지나 — 시간이 어디서 늘고 줄나","도입 제안 마무리",
 ["버전업에서 시간이 드는 일","지금 누가","크기","바뀌면"],
 [["병합·충돌 해결 + 리로케이션 재적용","사람","가장 큼",("LLM이 대신 → 줄어듦",GREEN)],
  ["관리 파일 ~25개 · 커밋 이력 갱신","사람","큼",("LLM이 초안 → 줄어듦",GREEN)],
  ["검사·테스트 실행","기계","—",("늘어도 사람 시간은 아님",GRAY)],
  ["확인(화면 열어보기·API 찔러보기)","사람","중간",("자동 검사가 대신 → 줄어듦",GREEN)],
  ["승인·판단","사람","작음",("그대로",GRAY)]],
 [4.6,1.4,1.2,4.9],BLUE,15,
 "첫 도입은 오히려 느리다(등록·기준선 세팅) / 두 번째 버전업부터 빨라진다.  보류 5종이 이 답을 좌우한다 — 사람 손이 붙는 것을 남기면 느려지고, 자동으로 도는 것만 남기면 빨라진다.")

# 15 안전장치
content("검사기의 정확한 위치 — LLM 위에 얹는 안전장치","오해 방지",
 [("검사기는 LLM을 대신하지 않는다. 검사도 판단도 LLM이 먼저 한다 — 검사기는 그 결과를 한 번 더 확인할 뿐이다",0,TEAL),
  ("1차 — LLM: 무엇을 바꿀지 정하고, 옮겨간 코드를 찾고, 충돌을 풀고, 검사까지 직접 돌린다. 속도는 전부 여기서 나온다",0,BLUE),
  ("2차 — 검사기: 같은 일을 다시 하지 않는다. LLM이 자기 자신에 대해 답할 수 없는 것만 확인한다",0,TEAL),
  ("① 정말 돌렸는가  ② 검사한 코드가 배포할 코드와 같은가  ③ 그 결과가 나중에 바뀌지 않았는가",1,DARK),
  ("그래서 독립성 — LLM이 실행·기록·증거·승인을 다 하면 ‘한 번 더’가 성립하지 않는다. 별도의 보호된 자동화(CI)가 실제 코드를 받아 검사를 돌려야 한다",0,RED),
  ("검사기를 ‘LLM보다 똑똑한 판정자’로 쓰면 실패한다 — 그건 이미 해봤다. ‘LLM이 한 일을 한 번 더 확인하는 안전장치’로 쓰면 부담은 작고 얻는 건 크다",0,GREEN)],TEAL,16)

# 16 결론
content("결론","한 장 요약",
 [("작은 변화 — 하위 패치엔 규칙 기반이 그대로 유효. 충돌 0·자동 병합 (4·5p)",0,GREEN),
  ("큰 변화 — 1.13.2처럼 구조가 바뀌면 오판(리로케이션) + 관리 파일 ~25개 운영 부담(6p) → 한계",0,RED),
  ("그래서 — 규칙 기반은 버리지 않는다. 예측 가능 범위엔 그대로, 불확실할 때만 하이브리드. 그 관리 산출물은 LLM이 만들고 검사기는 검증만 (10·11p)",0,BLUE),
  ("① 방식 통제 — 어떻게 일을 하는지를 지시(프롬프트·MD)로 통제한다 (순서·범위·남길 것)",0,TEAL),
  ("② 결과 통제 — 과정마다 정해진 형식으로 결과를 뱉게 강제 → 검사기가 판정",0,GREEN),
  ("한 줄 (harness = 고삐) — 업그레이드가 불확실해도 산출물을 예측 가능하게 만들어 통제한다. AI를 덜 쓰는(minimize) 게 아니라 고삐를 채우는(harness) 것",0,TEAL),
  ("현재: 1.13.1 9/9 · 1.13.2 부분 / 다음: 주장–증거 자동 대조기",2,GRAY)],TEAL,17)

prs.save("docs/CODEX-05-전략발표.pptx")
print("saved: docs/CODEX-05-전략발표.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
